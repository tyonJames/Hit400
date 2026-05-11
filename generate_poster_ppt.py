"""
Generate BlockLand Academic Poster as a single PowerPoint slide (A0 portrait).
Run: python generate_poster_ppt.py
Output: BlockLand_Poster.pptm
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os
import zipfile

# ── Palette (matches HIT template) ────────────────────────────────────────────
CREAM      = RGBColor(0xF5, 0xE5, 0xC8)   # warm background
ORANGE     = RGBColor(0xD4, 0xA0, 0x17)   # left-col headers (amber/gold)
BLUE       = RGBColor(0x1E, 0x3A, 0x8A)   # right-col headers (navy)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
DARK_RED   = RGBColor(0x8B, 0x0A, 0x14)   # top stripe
MID        = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xBB, 0xA0, 0x80)   # subtle cell border

# ── Layout constants (A0 portrait: 33.11" × 46.81") ──────────────────────────
SW = 33.11
SH = 46.81

M   = 0.30    # left/right outer margin
GAP = 0.18    # gap between columns
CW  = (SW - 2 * M - 3 * GAP) / 4    # column width ≈ 7.96"

# Column left edges
C1 = M
C2 = M + CW + GAP
C3 = M + 2 * (CW + GAP)
C4 = M + 3 * (CW + GAP)

HEADER_TOP  = 0.0
HEADER_H    = 4.00
CONTENT_TOP = 4.20
CONTENT_BOT = 45.50
FOOTER_TOP  = CONTENT_BOT + 0.10
FOOTER_H    = SH - FOOTER_TOP
CONTENT_H   = CONTENT_BOT - CONTENT_TOP   # ≈ 41.3"

SEC_HEADER_H = 0.70   # height of each coloured section header bar

DIAG_DIR = os.path.join(os.path.dirname(__file__), 'diagrams')

# ── Core helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w if line_w > Pt(0) else Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_txb(slide, text, l, t, w, h,
            size=20, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
            italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb, tf


def add_para(tf, text, size=20, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def embed_image(slide, path, l, t, w, h=None):
    """Add picture. If h is None, size by width only (preserve aspect ratio)."""
    if not os.path.exists(path):
        # Draw a placeholder box if image missing
        box = add_rect(slide, l, t, w, h or 3, fill=CREAM, line_color=BORDER, line_w=Pt(1))
        return box
    if h:
        pic = slide.shapes.add_picture(path, Inches(l), Inches(t),
                                       Inches(w), Inches(h))
    else:
        pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))
    return pic


# ── Section builder ───────────────────────────────────────────────────────────

def section(slide, x, y, w, total_h, title, header_color):
    """Draw coloured header bar + cream body. Returns body text frame."""
    # header bar
    hdr = add_rect(slide, x, y, w, SEC_HEADER_H, fill=header_color)
    hdr.line.fill.background()
    tf_h = hdr.text_frame
    tf_h.margin_left  = Inches(0.08)
    tf_h.margin_right = Inches(0.08)
    tf_h.margin_top   = Inches(0.06)
    p = tf_h.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # cream body
    body_h = total_h - SEC_HEADER_H
    body = add_rect(slide, x, y + SEC_HEADER_H, w, body_h,
                    fill=CREAM, line_color=BORDER, line_w=Pt(0.75))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.12)
    tf.margin_right  = Inches(0.12)
    tf.margin_top    = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    return tf, y + total_h   # text frame + next y position


# ── Poster builder ─────────────────────────────────────────────────────────────

def build_background(slide):
    add_rect(slide, 0, 0, SW, SH, fill=CREAM)


def build_header(slide):
    # Top colour stripe (dark red)
    add_rect(slide, 0, 0, SW, 0.25, fill=DARK_RED)

    # White header area
    add_rect(slide, 0, 0.25, SW, HEADER_H - 0.25, fill=WHITE)

    # HIT logo placeholder (left)
    logo_w, logo_h = 3.6, 3.0
    logo_box = add_rect(slide, M, 0.35, logo_w, logo_h,
                        fill=RGBColor(0xF0, 0xF0, 0xF0), line_color=BORDER, line_w=Pt(0.5))
    # Try to embed HIT logo if available
    logo_paths = ['diagrams/hit_logo.png', 'hit_logo.png']
    for lp in logo_paths:
        if os.path.exists(lp):
            slide.shapes.add_picture(lp, Inches(M), Inches(0.35), Inches(logo_w))
            break
    else:
        # Text fallback
        txb, tf = add_txb(slide, 'Harare Institute\nof Technology',
                           M + 0.1, 0.5, logo_w - 0.2, logo_h - 0.3,
                           size=28, bold=True, color=DARK_RED, align=PP_ALIGN.CENTER)

    # Project title (centre)
    title_l = M + logo_w + 0.3
    title_w = SW - title_l - M - 0.2
    txb, tf = add_txb(slide,
                      'BlockLand: A Blockchain-Based Land Registry and Management System for Zimbabwe',
                      title_l, 0.4, title_w, 2.0,
                      size=60, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Author / affiliation block
    txb, tf = add_txb(slide,
                      '1,2 Department of Information Security and Assurance, '
                      'School of Information Sciences and Technology, '
                      'Ganges Rd, Box BE 277, Belvedere, Harare, Zimbabwe',
                      title_l, 2.5, title_w, 0.65,
                      size=20, color=MID, align=PP_ALIGN.CENTER)

    txb, tf = add_txb(slide,
                      '1 Tinotenda James (H220349M),   2 Ms Chavhunduka',
                      title_l, 3.1, title_w, 0.45,
                      size=22, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    txb, tf = add_txb(slide,
                      '1 h220349m@hit.ac.zw,   2 chavhunduka@hit.ac.zw',
                      title_l, 3.55, title_w, 0.40,
                      size=18, italic=True, color=MID, align=PP_ALIGN.CENTER)

    # Bottom divider line under header
    add_rect(slide, 0, HEADER_H - 0.08, SW, 0.08, fill=DARK_RED)


def build_col1(slide):
    y = CONTENT_TOP
    body_sz = 19
    pad = 0.10

    # ── Abstract ──
    tf, y = section(slide, C1, y, CW, 9.5, 'Abstract', ORANGE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = (
        'BlockLand is a dual-ledger land registry system that anchors property '
        'ownership on the Stacks blockchain while maintaining a PostgreSQL database '
        'for operational metadata. The system eliminates paper-based title deed fraud '
        'through cryptographic ownership assertions, enforces a legally-aligned '
        'three-step transfer workflow via a Clarity smart contract, and provides a '
        'public verification portal accessible to any citizen without registration. '
        'JWT RS256 authentication, four-role RBAC (Admin, Registrar, Owner, Buyer), '
        'and IPFS document hashing provide defence-in-depth access control. Security '
        'testing confirmed resistance to replay attacks, privilege escalation, and '
        'double-transfer attempts. The system delivers sub-100 ms API response times '
        'at p95 under 50 concurrent users, with blockchain-confirming transfers '
        'completing in under 30 seconds on the Stacks testnet.'
    )

    y += pad

    # ── Introduction ──
    tf, y = section(slide, C1, y, CW, 8.5, 'Introduction', ORANGE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = (
        "Zimbabwe's Deeds Registry operates under the Deeds Registries Act "
        "[Chapter 20:05], maintaining physical title deeds and manual ledgers "
        "that are vulnerable to forgery, loss, and deliberate falsification. "
        "No public ownership verification mechanism exists — citizens, banks, "
        "and legal practitioners must submit formal inquiries to the Registry "
        "to confirm ownership. Land fraud costs Zimbabwean families their most "
        "significant asset with little recourse.\n\n"
        "Blockchain technology offers immutability, transparent auditability, "
        "and programmable multi-party workflows through smart contracts — "
        "properties directly applicable to land administration. BlockLand "
        "applies these capabilities to Zimbabwe's specific legal and "
        "administrative context without requiring a wholesale replacement "
        "of existing institutions."
    )

    y += pad

    # ── Problem Statement ──
    tf, y = section(slide, C1, y, CW, 12.5, 'Problem Statement', ORANGE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = (
        "Zimbabwe's land administration system relies on paper-based title deeds "
        "and a centralised manual ledger, creating five critical vulnerabilities:"
    )
    for item in [
        "Title deeds can be forged without any cryptographic reference.",
        "Ownership history is opaque — no tamper-evident audit trail exists.",
        "The centralised Deeds Registry is a single point of failure whose integrity "
        "cannot be independently verified by any external party.",
        "Property transfers require slow, paper-dependent in-person processes "
        "taking days to weeks.",
        "No public verification mechanism exists — ownership cannot be instantly "
        "confirmed by citizens, banks, or legal practitioners.",
    ]:
        pp = tf.add_paragraph()
        pp.alignment = PP_ALIGN.JUSTIFY
        pp.space_before = Pt(8)
        pp.level = 1
        run = pp.add_run()
        run.font.size = Pt(body_sz)
        run.font.color.rgb = BLACK
        run.text = f'•  {item}'

    y += pad

    # ── Objectives ──
    tf, y = section(slide, C1, y, CW, 10.3, 'Objectives', ORANGE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.text = 'This project aims to:'

    for num, obj in enumerate([
        'To provide a digital submission, review, and approval workflow for '
        'property title registrations — eliminating paper-based processes.',
        'To assign each property a unique on-chain token identifier, making it '
        'cryptographically impossible to register the same plot twice under '
        'different ownership records.',
        'To expose a public verification portal where any citizen can instantly '
        'confirm property ownership by plot number, deed number, or National ID — '
        'without requiring a registered account or government inquiry.',
    ], 1):
        pp = tf.add_paragraph()
        pp.alignment = PP_ALIGN.JUSTIFY
        pp.space_before = Pt(10)
        run = pp.add_run()
        run.font.size = Pt(body_sz)
        run.font.color.rgb = BLACK
        run.text = f'{num}.  {obj}'


def build_col2(slide):
    y = CONTENT_TOP
    pad = 0.10
    body_sz = 18

    # ── Use Cases ──
    tf, y = section(slide, C2, y, CW, 14.5, 'Use Cases', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.italic = True
    run.font.color.rgb = MID
    run.text = 'Key actor-system interactions'

    # Embed use case diagram
    uc_path = os.path.join(DIAG_DIR, '03_use_case.png')
    embed_image(slide,
                uc_path,
                C2 + 0.1,
                y - 14.5 + SEC_HEADER_H + 0.5,
                CW - 0.2,
                h=13.0)

    y += pad

    # ── System Overview ──
    tf, y = section(slide, C2, y, CW, 14.5, 'System Overview', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = (
        'BlockLand follows a three-tier architecture:\n\n'
        '   Presentation — Next.js 14 (React, TypeScript, TailwindCSS). '
        'Role-specific dashboards for Admin, Registrar, Owner, and Buyer. '
        'Public /verify portal requires no login.\n\n'
        '   Application — NestJS 10 REST API with JWT RS256 authentication '
        'and RBAC guards. BlockchainService adapts @stacks/transactions for '
        'all on-chain interactions.\n\n'
        '   Persistence — PostgreSQL 15 stores operational metadata. '
        'Stacks Blockchain (Clarity smart contract) stores immutable '
        'ownership assertions as canonical ground truth. IPFS stores '
        'title deed documents; content hash committed on-chain.\n\n'
        'The dual-ledger invariant: blockchain is authoritative; '
        'database is a queryable cache that must remain consistent '
        'with on-chain state.'
    )

    y += pad

    # ── Transfer Workflow (UML) ──
    tf, y = section(slide, C2, y, CW, 12.3, 'Use Cases / Transfer Workflow (UML)', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.italic = True
    run.font.color.rgb = MID
    run.text = 'Three-step on-chain transfer workflow'

    flow_path = os.path.join(DIAG_DIR, '07_activity.png')
    embed_image(slide,
                flow_path,
                C2 + 0.1,
                y - 12.3 + SEC_HEADER_H + 0.5,
                CW - 0.2,
                h=11.0)


def build_col3(slide):
    y = CONTENT_TOP
    pad = 0.10
    body_sz = 18

    # ── Architecture Diagram (large) ──
    arch_sec_h = 22.5
    tf, y_after = section(slide, C3, y, CW, arch_sec_h, 'Architecture Diagram', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.italic = True
    run.font.color.rgb = MID
    run.text = 'Three-tier solution architecture'

    arch_path = os.path.join(DIAG_DIR, '06_architecture.png')
    embed_image(slide,
                arch_path,
                C3 + 0.05,
                y + SEC_HEADER_H + 0.4,
                CW - 0.1,
                h=arch_sec_h - SEC_HEADER_H - 0.6)

    y = y_after + pad

    # ── Results and Screenshots ──
    tf, y = section(slide, C3, y, CW, 18.8, 'Results and Screenshots', ORANGE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.text = 'Security Test Results (all adversarial tests PASS):'

    tests = [
        ('JWT replay after 15-min expiry',       'PASS — 401 Unauthorized'),
        ('JWT HS256 downgrade forgery attempt',  'PASS — 401 Unauthorized'),
        ('Horizontal privilege escalation',       'PASS — 403 Forbidden'),
        ('Double-transfer same plot',             'PASS — (err u106)'),
        ('Direct contract call, non-registrar',  'PASS — (err u101)'),
        ('IPFS hash substitution via contract',  'PASS — (err u101)'),
    ]
    for test, result in tests:
        pp = tf.add_paragraph()
        pp.space_before = Pt(6)
        run = pp.add_run()
        run.font.size = Pt(17)
        run.font.color.rgb = BLACK
        run.text = f'  {test}'
        run2 = pp.add_run()
        run2.font.size = Pt(17)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0x15, 0x6E, 0x30)
        run2.text = f'  ✔ {result}'

    pp = tf.add_paragraph()
    pp.space_before = Pt(14)
    run = pp.add_run()
    run.font.size = Pt(body_sz)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.text = 'Performance Metrics (50 concurrent users):'

    for label, val in [
        ('Property listing (GET /properties)', 'p95 = 45 ms'),
        ('Public verify portal (/verify/:plot)', 'p95 = 28 ms'),
        ('Blockchain transfer (initiate)',       'p95 = 2,100 ms'),
        ('End-to-end transfer (testnet)',        '< 30 seconds'),
    ]:
        pp = tf.add_paragraph()
        pp.space_before = Pt(5)
        run = pp.add_run()
        run.font.size = Pt(17)
        run.font.color.rgb = BLACK
        run.text = f'  {label}:  '
        run2 = pp.add_run()
        run2.font.size = Pt(17)
        run2.font.bold = True
        run2.font.color.rgb = BLUE
        run2.text = val

    pp = tf.add_paragraph()
    pp.space_before = Pt(14)
    run = pp.add_run()
    run.font.size = Pt(body_sz)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.text = 'Functional Verification:'

    for item in [
        'All 7 FRs derived from Deeds Registries Act verified via Clarinet test suite.',
        'Three-step transfer workflow confirmed end-to-end on Stacks testnet.',
        'Duplicate plot registration rejected at contract level (ERR-TITLE-DEED-EXISTS).',
        'Public /verify portal returns ownership status in < 50 ms, no login required.',
    ]:
        pp = tf.add_paragraph()
        pp.space_before = Pt(6)
        run = pp.add_run()
        run.font.size = Pt(17)
        run.font.color.rgb = BLACK
        run.text = f'  •  {item}'

    # Embed interface screenshot
    iface_path = os.path.join(DIAG_DIR, '12_interface.png')
    embed_image(slide,
                iface_path,
                C3 + 0.1,
                y - 3.2,
                CW - 0.2,
                h=3.0)


def build_col4(slide):
    y = CONTENT_TOP
    pad = 0.10
    body_sz = 18

    # ── Conclusion ──
    tf, y = section(slide, C4, y, CW, 14.5, 'Conclusion', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = (
        'BlockLand demonstrates that targeted blockchain integration — '
        'anchoring only ownership assertions on-chain while preserving '
        'relational databases for operational richness — delivers meaningful '
        'security guarantees for Zimbabwe\'s land administration context.\n\n'
        'The Clarity smart contract enforces the legally-required '
        'three-party transfer consent model (seller, buyer, registrar) '
        'as required by the Deeds Registries Act [Chapter 20:05], with '
        'independent signature verification at each step. No single party '
        'can unilaterally reassign ownership.\n\n'
        'IPFS document hashing commits title deed content on-chain, '
        'making silent document substitution cryptographically detectable. '
        'The title-deed-index map makes duplicate plot registration '
        'impossible.\n\n'
        'All seven principal functional requirements and six adversarial '
        'security test cases are verified. The system achieves sub-100 ms '
        'response times for non-blockchain operations, demonstrating '
        'practical deployability within Zimbabwe\'s infrastructure context.'
    )

    y += pad

    # ── Future Work ──
    tf, y = section(slide, C4, y, CW, 9.5, 'Future Work', BLUE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.font.size = Pt(body_sz)
    run.font.color.rgb = BLACK
    run.text = 'Planned improvements:'

    for item in [
        'Client-side Hiro Wallet signing — individual users sign their own '
        'blockchain transactions for non-repudiable per-user records.',
        'Stacks mainnet deployment — transfer full Bitcoin security '
        'guarantees to the ownership record.',
        'Zero-knowledge identity proofs — replace National ID transmission '
        'with ZK proofs to protect citizen privacy.',
        'Legislative engagement — support amendment of the Electronic '
        'Transactions and Electronic Commerce Act to legally recognise '
        'blockchain land records.',
        'National ID and owner name search on the public /verify portal '
        'to replace the current internal UUID-based lookup.',
    ]:
        pp = tf.add_paragraph()
        pp.space_before = Pt(8)
        run = pp.add_run()
        run.font.size = Pt(body_sz - 1)
        run.font.color.rgb = BLACK
        run.text = f'•  {item}'

    y += pad

    # ── References ──
    tf, y = section(slide, C4, y, CW, 17.3, 'References', BLUE)
    refs = [
        'M. Zein, R. & Twinomurinzi, H. (2023). Blockchain Technology in Lands '
        'Registration: A Systematic Literature Review. JeDEM, 15(2), pp. 1-36.',
        'Mansoor, M.A. et al. (2023). Blockchain Technology for Land Registry '
        'Management in Developing Countries. ETECTE Conf., Lahore.',
        'Alam, K.M. et al. (2022). A Blockchain-based Land Title Management System '
        'for Bangladesh. J. King Saud Univ., 34(6), pp. 3096-3110.',
        'Paavo, J.P. et al. (2025). Practicality of Blockchain for Land Registration: '
        'A Namibian Case Study. Land, 14(8), p. 1626.',
        'Konashevych, O. (2020). General Concept of Real Estate Tokenization on '
        'Blockchain. European Property Law Journal, 9(1), pp. 21-66.',
        'Dwivedi, R. et al. (2023). Blockchain-Based Transferable Digital Rights '
        'of Land. arXiv:2308.05950.',
        'Kombe, C. et al. (2017). Design of Land Administration and Title '
        'Registration Model Based on Blockchain Technology.',
        'Kimathi, K.T. (n.d.). A Land Administration Model Based on Blockchain '
        'Technology. [Unpublished thesis]. HIT, Harare.',
        'Ameyaw, P.D. & de Vries, W.T. (2023). Blockchain technology adaptation '
        'for land administration: socio-cultural elements. Land Use Policy.',
        'Nakamoto, S. (2008). Bitcoin: A Peer-to-Peer Electronic Cash System.',
    ]
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.size = Pt(16)
    run.font.color.rgb = BLACK
    run.text = refs[0]

    for ref in refs[1:]:
        pp = tf.add_paragraph()
        pp.space_before = Pt(6)
        run = pp.add_run()
        run.font.size = Pt(16)
        run.font.color.rgb = BLACK
        run.text = ref


def build_footer(slide):
    # Dark red footer bar
    add_rect(slide, 0, FOOTER_TOP, SW, 0.12, fill=DARK_RED)

    foot_bg_h = SH - FOOTER_TOP - 0.12
    add_rect(slide, 0, FOOTER_TOP + 0.12, SW, foot_bg_h, fill=WHITE)

    # Year (large, left)
    txb, _ = add_txb(slide, '2025',
                     M, FOOTER_TOP + 0.15, 2.5, 1.1,
                     size=60, bold=True, color=BLACK)

    # HIT address (centre)
    txb, _ = add_txb(slide,
                     'Harare Institute of Technology\n'
                     'P.O. Box BE 277, Belvedere, Harare, Zimbabwe\n'
                     'success through innovation',
                     M + 2.8, FOOTER_TOP + 0.18, SW - M - 2.8 - 0.5, 1.0,
                     size=18, color=MID, align=PP_ALIGN.CENTER, italic=False)

    # Bottom stripe
    add_rect(slide, 0, SH - 0.12, SW, 0.12, fill=DARK_RED)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    build_background(slide)
    build_header(slide)
    build_col1(slide)
    build_col2(slide)
    build_col3(slide)
    build_col4(slide)
    build_footer(slide)

    tmp    = 'BlockLand_Poster_tmp.pptx'
    output = 'BlockLand_Poster.pptm'

    prs.save(tmp)

    # Rewrite as .pptm by patching the content-type in the ZIP
    pptx_ct = (b'application/vnd.openxmlformats-officedocument'
               b'.presentationml.presentation.main+xml')
    pptm_ct = (b'application/vnd.ms-powerpoint'
               b'.presentation.macroEnabled.main+xml')

    with zipfile.ZipFile(tmp, 'r') as zin, \
         zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == '[Content_Types].xml':
                data = data.replace(pptx_ct, pptm_ct)
            zout.writestr(item, data)

    os.remove(tmp)
    print(f'Saved: {output}')

if __name__ == '__main__':
    main()
