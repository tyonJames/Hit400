"""
Generate BlockLand HIT capstone presentation matching the official template.
Run: python generate_hit_ppt.py
Output: BlockLand_HIT_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette (matched to template) ─────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x0D, 0x17, 0x41)   # title bands, left footer
TEAL       = RGBColor(0x1A, 0x6B, 0x55)   # top bar, right footer
GOLD       = RGBColor(0xB8, 0x86, 0x0B)   # separator line, slide numbers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)

# ── Slide dimensions (standard 16:9) ──────────────────────────────────────────
SW = 13.33
SH = 7.50

DEPT       = 'DEPARTMENT OF INFORMATION SECURITY AND ASSURANCE'
RUNNING    = 'BlockLand: Blockchain-Based Land Registry System for Zimbabwe'
DIAG_DIR   = os.path.join(os.path.dirname(__file__), 'diagrams')

# ── Core helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w if line_w > Pt(0) else Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_txb(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, italic=False, wrap=True, font='Calibri'):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb, tf


def add_para(tf, text, size=20, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=8, italic=False, font='Calibri', level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def embed_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return add_rect(slide, l, t, w, h or 3, fill=LIGHT_GRAY,
                        line_color=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(1))
    if h:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))


# ── Chrome builders ───────────────────────────────────────────────────────────

def chrome_content(slide, slide_title, slide_num):
    """Standard chrome for all content slides (slides 2+)."""
    # ── Thin teal top bar ────────────────────────────────────────────────────
    add_rect(slide, 0, 0, SW, 0.30, fill=TEAL)
    add_txb(slide, RUNNING, 0.1, 0, SW - 0.2, 0.30,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Dark navy title band ─────────────────────────────────────────────────
    add_rect(slide, 0, 0.30, SW, 1.10, fill=DARK_NAVY)
    add_txb(slide, slide_title, 0.30, 0.32, SW - 0.5, 1.06,
            size=40, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
            font='Calibri Light')

    # ── Gold separator line ──────────────────────────────────────────────────
    add_rect(slide, 0, 1.40, SW, 0.045, fill=GOLD)

    # ── Footer bar ───────────────────────────────────────────────────────────
    footer_y = 6.95
    add_rect(slide, 0,        footer_y, 7.20, 0.55, fill=DARK_NAVY)
    add_rect(slide, 7.20,     footer_y, 5.63, 0.55, fill=TEAL)
    add_rect(slide, SW - 0.5, footer_y, 0.50, 0.55, fill=GOLD)

    add_txb(slide, DEPT,
            0.15, footer_y + 0.08, 7.0, 0.40,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_txb(slide, 'HARARE INSTITUTE OF TECHNOLOGY',
            7.30, footer_y + 0.08, 5.4, 0.40,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_txb(slide, str(slide_num),
            SW - 0.48, footer_y + 0.06, 0.44, 0.43,
            size=16, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # ── Small HIT logo placeholder (bottom-right of content area) ────────────
    logo_paths = ['diagrams/hit_logo.png', 'hit_logo.png']
    for lp in logo_paths:
        if os.path.exists(lp):
            slide.shapes.add_picture(lp, Inches(11.55), Inches(5.60), Inches(1.55))
            break
    else:
        # Text fallback
        add_txb(slide, 'HIT', 12.0, 5.75, 1.1, 0.6,
                size=22, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
        add_txb(slide, 'Harare Institute\nof Technology',
                11.55, 6.10, 1.65, 0.75,
                size=7, color=DARK_NAVY, align=PP_ALIGN.CENTER)


def content_tf(slide, l=0.45, t=1.55, w=10.8, h=5.25):
    """Return a text frame in the content area."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    return tf


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top navy band
    add_rect(slide, 0, 0, SW, 1.75, fill=DARK_NAVY)

    # "TITLE OF PROJECT" in navy band
    title_txt = ('BlockLand: A Blockchain-Based Land Registry\n'
                 'and Management System for Zimbabwe')
    add_txb(slide, title_txt,
            0.4, 0.08, SW - 0.8, 1.60,
            size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            font='Calibri')

    # Thin teal stripe below navy
    add_rect(slide, 0, 1.75, SW, 0.10, fill=TEAL)

    # Gold line
    add_rect(slide, 0, 1.85, SW, 0.045, fill=GOLD)

    # White content area
    cx, cy, cw = 0.9, 2.10, SW - 1.8

    # Thesis subtitle
    add_txb(slide,
            'A capstone thesis submitted in partial fulfilment of the requirements for the degree of',
            cx, cy, cw, 0.40,
            size=13, italic=True, color=BLACK, align=PP_ALIGN.CENTER)

    add_txb(slide,
            'BACHELOR OF TECHNOLOGY in INFORMATION SECURITY AND ASSURANCE',
            cx, cy + 0.38, cw, 0.38,
            size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    add_txb(slide, 'by', cx, cy + 0.78, cw, 0.28,
            size=13, italic=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Student name
    add_txb(slide, 'Tinotenda James', cx, cy + 1.05, cw, 0.48,
            size=28, bold=True, color=BLACK, align=PP_ALIGN.CENTER,
            font='Palatino Linotype')

    add_txb(slide, 'H220349M', cx, cy + 1.52, cw, 0.28,
            size=13, color=BLACK, align=PP_ALIGN.CENTER)
    add_txb(slide, 'h220349m@hit.ac.zw', cx, cy + 1.78, cw, 0.28,
            size=13, color=BLACK, align=PP_ALIGN.CENTER)

    # Supervisor
    add_txb(slide, 'Under the guidance of', cx, cy + 2.15, cw, 0.28,
            size=13, italic=True, color=BLACK, align=PP_ALIGN.CENTER)

    add_txb(slide, 'Ms Chavhunduka', cx, cy + 2.42, cw, 0.48,
            size=24, bold=True, color=BLACK, align=PP_ALIGN.CENTER,
            font='Palatino Linotype')

    add_txb(slide, 'LECTURER', cx, cy + 2.88, cw, 0.28,
            size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Institution block
    add_txb(slide, 'School of Information Sciences and Technology',
            cx, cy + 3.30, cw, 0.30,
            size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_txb(slide, 'Harare Institute of Technology',
            cx, cy + 3.58, cw, 0.28,
            size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_txb(slide, 'BE277, Ganges Road, Belvedere, Harare',
            cx, cy + 3.84, cw, 0.26,
            size=11, color=BLACK, align=PP_ALIGN.CENTER)
    add_txb(slide, 'May 2026',
            cx, cy + 4.10, cw, 0.30,
            size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Footer
    footer_y = 6.95
    add_rect(slide, 0,        footer_y, 7.20, 0.55, fill=DARK_NAVY)
    add_rect(slide, 7.20,     footer_y, 5.63, 0.55, fill=TEAL)
    add_rect(slide, SW - 0.5, footer_y, 0.50, 0.55, fill=GOLD)
    add_txb(slide, DEPT,
            0.15, footer_y + 0.08, 7.0, 0.40,
            size=9, bold=True, color=WHITE)
    add_txb(slide, 'HARARE INSTITUTE OF TECHNOLOGY',
            7.30, footer_y + 0.08, 5.4, 0.40,
            size=9, bold=True, color=WHITE)
    add_txb(slide, '1', SW - 0.48, footer_y + 0.06, 0.44, 0.43,
            size=16, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # HIT logo (right of supervisor)
    logo_paths = ['diagrams/hit_logo.png', 'hit_logo.png']
    for lp in logo_paths:
        if os.path.exists(lp):
            slide.shapes.add_picture(lp, Inches(5.9), Inches(5.2), Inches(1.6))
            break


# ── Slide 2 — Overview ────────────────────────────────────────────────────────

def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Overview', 2)
    tf = content_tf(slide)

    p = tf.paragraphs[0]
    p.space_before = Pt(4)

    for item in [
        'Introduction',
        'Problem Statement',
        'Objectives',
        'Related Works',
        'Proposed Solution',
        'System Design & Smart Contract',
        'Security Analysis',
        'Results & Evaluation',
        'Conclusion and Future Work',
        'References',
    ]:
        if item == 'Introduction':
            run = p.add_run()
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)

        run = p.add_run()
        run.text = f'❑  {item}'     # ❑ checkbox symbol
        run.font.name = 'Calibri'
        run.font.size = Pt(24)
        run.font.color.rgb = BLACK


# ── Slide 3 — Introduction ────────────────────────────────────────────────────

def slide_introduction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Introduction', 3)
    tf = content_tf(slide)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Land Administration in Zimbabwe'
    run.font.name = 'Calibri'
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    for bullet in [
        'Land is the most significant asset class in Zimbabwe, yet the Deeds Registry '
        'operates under a paper-based, centralised system (Deeds Registries Act [Chapter 20:05]).',
        'Physical title deeds are vulnerable to forgery, loss, and manipulation with no '
        'cryptographic proof of authenticity.',
        'No public mechanism exists for citizens, banks, or lawyers to independently '
        'verify property ownership without a formal Registry inquiry.',
    ]:
        add_para(tf, f'•  {bullet}', size=19, space_before=10, color=BLACK)

    add_para(tf, 'Why Blockchain?', size=22, bold=True, color=DARK_NAVY, space_before=18)

    for bullet in [
        'Immutability — ownership records cannot be silently altered once committed on-chain.',
        'Transparency — any party can audit the full chain of ownership independently.',
        'Programmability — smart contracts enforce multi-party transfer workflows, '
        'eliminating procedural fraud.',
        'Bitcoin-anchored security — BlockLand uses the Stacks blockchain, which settles '
        'each block to Bitcoin via Proof of Transfer (PoX).',
    ]:
        add_para(tf, f'•  {bullet}', size=19, space_before=8, color=BLACK)


# ── Slide 4 — Problem Statement ───────────────────────────────────────────────

def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Problem Statement', 4)
    tf = content_tf(slide)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Zimbabwe's land administration system relies on paper-based title deeds and a "
                "centralised manual ledger, creating five critical vulnerabilities:")
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.color.rgb = DARK_NAVY

    problems = [
        ('Title Deed Forgery',
         'Physical documents can be duplicated or falsified with no cryptographic reference '
         'for verification.'),
        ('Opacity of Ownership History',
         'No publicly accessible, tamper-evident record of who owned a property and when.'),
        ('Single Point of Failure',
         'The centralised Deeds Registry is the sole authority — its integrity cannot be '
         'independently verified by any external party.'),
        ('Slow Transfer Process',
         'Property transfers require physical document exchange and in-person approvals, '
         'taking days to weeks.'),
        ('No Public Verification',
         'Citizens, banks, and legal practitioners cannot instantly confirm ownership without '
         'a formal inquiry to the Registry.'),
    ]

    for i, (title, desc) in enumerate(problems, 1):
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f'{i}.  {title}: '
        run.font.name = 'Calibri'
        run.font.size = Pt(19)
        run.font.bold = True
        run.font.color.rgb = GOLD
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(19)
        run2.font.bold = False
        run2.font.color.rgb = BLACK


# ── Slide 5 — Objectives ──────────────────────────────────────────────────────

def slide_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Objectives', 5)
    tf = content_tf(slide, t=1.60, h=5.10)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'This project aims:'
    run.font.name = 'Calibri'
    run.font.size = Pt(21)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    objectives = [
        ('Objective 1 — Digital Registration Workflow',
         'To provide a digital submission, review, and approval workflow for property title '
         'registrations — eliminating paper-based processes and reducing registration time '
         'from weeks to minutes.'),
        ('Objective 2 — On-Chain Unique Token ID',
         'To assign each property a unique on-chain token identifier, making it '
         'cryptographically impossible to register the same plot twice under different '
         'ownership records.'),
        ('Objective 3 — Public Verification Portal',
         'To expose a public verification portal where any citizen can instantly confirm '
         'property ownership by plot number, deed number, or National ID — without '
         'requiring a registered account or government inquiry.'),
    ]

    for title, desc in objectives:
        add_para(tf, '', size=12, space_before=14)
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = f'✔  {title}'
        run.font.name = 'Calibri'
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TEAL

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        # Set left margin via XML (python-pptx 1.x removed paragraph_format)
        pPr = p2._p.get_or_add_pPr()
        pPr.set('marL', str(int(914400 * 0.35)))
        run2 = p2.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK


# ── Slide 6 — Related Works ───────────────────────────────────────────────────

def slide_related(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Related Works', 6)
    tf = content_tf(slide, h=5.20)

    works = [
        ('Kombe et al. (2017)',
         'First formal model proposing a three-party workflow (seller, buyer, registrar) '
         'for blockchain land registration — the pattern BlockLand enforces on-chain.'),
        ('Alam et al. (2022)',
         'Ethereum-based land title system for Bangladesh. Demonstrated feasibility but '
         'highlighted gas cost as a prohibitive barrier — absent on Stacks.'),
        ('M. Zein & Twinomurinzi (2023)',
         'Systematic review of 47 papers — fewer than a third provided any formal threat '
         'modelling. Identified security analysis as the key gap in the literature.'),
        ('Mansoor et al. (2023)',
         'Recommended hybrid architectures (blockchain + relational database) as the most '
         'practically deployable approach for developing economy contexts.'),
        ('Paavo et al. (2025)',
         'Namibian pilot study — identified registrar workflow integration as the binding '
         'success factor, directly motivating BlockLand\'s registrar finalisation step.'),
        ('Konashevych (2020)',
         'Identified the oracle problem — gap between on-chain digital state and physical '
         'reality — as the most fundamental limitation of blockchain in land administration.'),
    ]

    p = tf.paragraphs[0]
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = works[0][0] + ':  '
    run.font.name = 'Calibri'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY
    run2 = p.add_run()
    run2.text = works[0][1]
    run2.font.name = 'Calibri'
    run2.font.size = Pt(18)
    run2.font.color.rgb = BLACK

    for author, desc in works[1:]:
        p = tf.add_paragraph()
        p.space_before = Pt(9)
        run = p.add_run()
        run.text = f'{author}:  '
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = DARK_NAVY
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK


# ── Slide 7 — Proposed Solution ───────────────────────────────────────────────

def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Proposed Solution', 7)

    # Left text column
    tf = content_tf(slide, l=0.45, t=1.58, w=5.8, h=5.20)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Three-Tier Architecture'
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    tiers = [
        ('Presentation', 'Next.js 14 — Role-specific dashboards, Hiro Wallet integration, '
         'public /verify portal.'),
        ('Application', 'NestJS 10 REST API — JWT RS256 auth, RBAC guards, '
         'BlockchainService adapter.'),
        ('Persistence', 'PostgreSQL — operational metadata.  '
         'Stacks Blockchain — immutable ownership ground truth via Clarity smart contract.'),
    ]
    for tier, desc in tiers:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f'▶  {tier}: '
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK

    add_para(tf, 'Dual-Ledger Invariant', size=20, bold=True,
             color=DARK_NAVY, space_before=18)
    add_para(tf,
             'Blockchain = canonical ownership ground truth.  '
             'PostgreSQL = queryable operational cache.  '
             'If the two ever diverge, the blockchain record takes precedence.',
             size=17, color=BLACK, space_before=6)

    add_para(tf, 'IPFS Document Anchoring', size=20, bold=True,
             color=DARK_NAVY, space_before=14)
    add_para(tf,
             'Title deed documents stored on IPFS. Content hash committed on-chain — '
             'any document substitution is cryptographically detectable.',
             size=17, color=BLACK, space_before=6)

    # Right: architecture diagram
    arch_path = os.path.join(DIAG_DIR, '06_architecture.png')
    embed_image(slide, arch_path, l=6.50, t=1.58, w=6.60, h=5.20)


# ── Slide 8 — System Design & Smart Contract ──────────────────────────────────

def slide_smart_contract(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'System Design & Smart Contract', 8)

    # Left column: contract maps
    tf = content_tf(slide, l=0.45, t=1.58, w=6.0, h=5.20)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'blockland.clar — On-Chain Data Maps'
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    maps = [
        ('property-registry',       'uint → {owner, title-deed-hash, status, ipfs-hash}'),
        ('title-deed-index',        '(buff 32) → uint   [fraud prevention reverse lookup]'),
        ('transfer-requests',       'uint → {seller, buyer, buyer-approved, registrar-approved}'),
        ('ownership-history',       '{property-id, seq} → {owner, acquired-at}'),
        ('authorized-registrars',   'principal → bool   [registrar whitelist]'),
    ]
    for name, desc in maps:
        p = tf.add_paragraph()
        p.space_before = Pt(9)
        run = p.add_run()
        run.text = f'  {name}'
        run.font.name = 'Courier New'
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run2 = p.add_run()
        run2.text = f'\n     {desc}'
        run2.font.name = 'Courier New'
        run2.font.size = Pt(13)
        run2.font.color.rgb = BLACK

    add_para(tf, 'Three-Step Transfer Workflow', size=20, bold=True,
             color=DARK_NAVY, space_before=16)

    steps = [
        ('Step 1 — Owner', 'Calls initiate-transfer — property locked to "pending-transfer"'),
        ('Step 2 — Buyer', 'Calls buyer-approve-transfer — sets buyer-approved = true'),
        ('Step 3 — Registrar', 'Calls registrar-finalize-transfer — ownership atomically reassigned'),
    ]
    for label, desc in steps:
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f'✔  {label}: '
        run.font.name = 'Calibri'
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = GOLD
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(17)
        run2.font.color.rgb = BLACK

    # Right: use case / activity diagram
    act_path = os.path.join(DIAG_DIR, '07_activity.png')
    embed_image(slide, act_path, l=6.65, t=1.58, w=6.45, h=5.20)


# ── Slide 9 — Security Analysis ───────────────────────────────────────────────

def slide_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Security Analysis', 9)
    tf = content_tf(slide, h=5.20)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'STRIDE Threat Model — Applied to BlockLand'
    run.font.name = 'Calibri'
    run.font.size = Pt(21)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    stride = [
        ('S — Spoofing',
         'JWT RS256 asymmetric signing + on-chain registrar whitelist prevent identity impersonation.'),
        ('T — Tampering',
         'Blockchain is canonical — DB inconsistency detectable by querying contract state. '
         'IPFS hash prevents document substitution.'),
        ('R — Repudiation',
         'Every transfer step is a signed Stacks transaction — immutable on-chain record '
         'with block height timestamp.'),
        ('I — Information Disclosure',
         'JWT-gated API; RBAC restricts financial data to parties in each transaction.'),
        ('D — Denial of Service',
         'Rate limiting at API gateway; blockchain operates independently of API availability.'),
        ('E — Elevation of Privilege',
         'Dual-layer enforcement — NestJS RolesGuard AND Clarity contract tx-sender check. '
         'Both must be bypassed simultaneously for escalation to succeed.'),
    ]

    for category, mitigation in stride:
        p = tf.add_paragraph()
        p.space_before = Pt(9)
        run = p.add_run()
        run.text = f'{category}:  '
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run2 = p.add_run()
        run2.text = mitigation
        run2.font.name = 'Calibri'
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK

    add_para(tf,
             'Clarity language eliminates reentrancy and unbounded-loop vulnerabilities '
             'by design — all execution paths are statically analysable before deployment.',
             size=17, italic=True, color=DARK_NAVY, space_before=14)


# ── Slide 10 — Results & Evaluation ──────────────────────────────────────────

def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Results & Evaluation', 10)

    # Left column: security tests
    tf_l = content_tf(slide, l=0.45, t=1.58, w=6.2, h=5.20)

    p = tf_l.paragraphs[0]
    run = p.add_run()
    run.text = 'Adversarial Security Tests'
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    tests = [
        ('JWT replay after 15-min expiry',      'PASS — 401'),
        ('JWT HS256 downgrade forgery',          'PASS — 401'),
        ('Horizontal privilege escalation',       'PASS — 403'),
        ('Double-transfer same plot',             'PASS — err u106'),
        ('Direct contract call (non-registrar)', 'PASS — err u101'),
        ('IPFS hash substitution via contract',  'PASS — err u101'),
    ]
    for test, result in tests:
        p = tf_l.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f'•  {test}  '
        run.font.name = 'Calibri'
        run.font.size = Pt(17)
        run.font.color.rgb = BLACK
        run2 = p.add_run()
        run2.text = result
        run2.font.name = 'Calibri'
        run2.font.size = Pt(17)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0x15, 0x6E, 0x30)

    add_para(tf_l, 'Functional Requirements', size=20, bold=True,
             color=DARK_NAVY, space_before=16)
    add_para(tf_l,
             'All 7 FRs derived from the Deeds Registries Act verified via '
             'Clarinet unit tests and adversarial integration testing.',
             size=17, color=BLACK, space_before=6)

    # Right column: performance
    tf_r = content_tf(slide, l=6.80, t=1.58, w=6.30, h=5.20)

    p = tf_r.paragraphs[0]
    run = p.add_run()
    run.text = 'Performance Metrics  (50 concurrent users)'
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    metrics = [
        ('GET /properties',         'p95 = 45 ms'),
        ('Public /verify portal',   'p95 = 28 ms'),
        ('POST /properties',        'p95 = 89 ms'),
        ('Blockchain initiate-tx',  'p95 = 2,100 ms'),
        ('End-to-end transfer',     '< 30 seconds'),
    ]
    for label, val in metrics:
        p = tf_r.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f'{label}: '
        run.font.name = 'Calibri'
        run.font.size = Pt(19)
        run.font.color.rgb = BLACK
        run2 = p.add_run()
        run2.text = val
        run2.font.name = 'Calibri'
        run2.font.size = Pt(19)
        run2.font.bold = True
        run2.font.color.rgb = TEAL

    add_para(tf_r, 'Key Findings', size=20, bold=True, color=DARK_NAVY, space_before=20)

    for finding in [
        'title-deed-index map makes duplicate plot registration cryptographically impossible.',
        'Clarity\'s non-Turing-complete design eliminates reentrancy by language specification.',
        'Three-step workflow confirmed end-to-end on Stacks testnet — no unilateral transfer possible.',
        'Public /verify portal returns ownership status in < 50 ms with no login required.',
    ]:
        add_para(tf_r, f'✔  {finding}', size=17, color=BLACK, space_before=8)


# ── Slide 11 — Conclusion & Future Works ──────────────────────────────────────

def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'Conclusion and Future Works', 11)

    # Left: conclusion
    tf_l = content_tf(slide, l=0.45, t=1.58, w=6.2, h=5.20)

    p = tf_l.paragraphs[0]
    run = p.add_run()
    run.text = 'Conclusion'
    run.font.name = 'Calibri'
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    conclusions = [
        'BlockLand demonstrates that targeted blockchain integration — anchoring only '
        'ownership assertions on-chain — delivers meaningful security guarantees for '
        'Zimbabwe\'s land administration context.',
        'The Clarity smart contract enforces the three-party consent model required by '
        'the Deeds Registries Act [Chapter 20:05] at the blockchain layer — no application '
        'code can bypass it.',
        'IPFS document hashing and the title-deed-index reverse lookup map prevent '
        'document substitution and duplicate plot registration respectively.',
        'All 7 functional requirements and 6 adversarial security tests verified.',
    ]
    for c in conclusions:
        add_para(tf_l, f'•  {c}', size=18, color=BLACK, space_before=10)

    # Right: future work
    tf_r = content_tf(slide, l=6.80, t=1.58, w=6.30, h=5.20)

    p = tf_r.paragraphs[0]
    run = p.add_run()
    run.text = 'Future Work'
    run.font.name = 'Calibri'
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY

    futures = [
        ('Client-Side Wallet Signing',
         'Integrate Hiro Wallet openContractCall() so individual users sign their own '
         'blockchain transactions — creating non-repudiable per-user records.'),
        ('Stacks Mainnet Deployment',
         'Deploy to Stacks mainnet, transferring the full Bitcoin security guarantee to '
         'the ownership record.'),
        ('Zero-Knowledge Identity Proofs',
         'Replace National ID transmission with ZK proofs — verifiable without exposing '
         'the citizen\'s raw identity data.'),
        ('Legislative Engagement',
         'Support amendment of the Electronic Transactions Act to formally recognise '
         'blockchain land records as legally binding.'),
    ]
    for title, desc in futures:
        p = tf_r.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = f'▶  {title}: '
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = TEAL
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = 'Calibri'
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK


# ── Slide 12 — References ─────────────────────────────────────────────────────

def slide_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chrome_content(slide, 'References', 12)
    tf = content_tf(slide, t=1.55, h=5.25)

    refs = [
        'Alam, K.M. et al. (2022). A Blockchain-based Land Title Management System for Bangladesh. '
        'J. King Saud University — Computer and Information Sciences, 34(6), pp. 3096-3110.',
        'Ameyaw, P.D. & de Vries, W.T. (2023). Blockchain technology adaptation for land '
        'administration services: The importance of socio-cultural elements. Land Use Policy.',
        'Dwivedi, R. et al. (2023). Blockchain-Based Transferable Digital Rights of Land. '
        'arXiv:2308.05950.',
        'Kimathi, K.T. (n.d.). A Land Administration Model Based on Blockchain Technology. '
        '[Unpublished thesis]. Harare Institute of Technology.',
        'Kombe, C. et al. (2017). Design of Land Administration and Title Registration Model '
        'Based on Blockchain Technology.',
        'Konashevych, O. (2020). General Concept of Real Estate Tokenization on Blockchain: '
        'The Right to Choose. European Property Law Journal, 9(1), pp. 21-66.',
        'M. Zein, R. & Twinomurinzi, H. (2023). Blockchain Technology in Lands Registration: '
        'A Systematic Literature Review. JeDEM, 15(2), pp. 1-36.',
        'Mansoor, M.A. et al. (2023). Blockchain Technology for Land Registry Management in '
        'Developing Countries. 2023 ETECTE Conference, Lahore.',
        'Nakamoto, S. (2008). Bitcoin: A Peer-to-Peer Electronic Cash System.',
        'Paavo, J.P. et al. (2025). Practicality of Blockchain Technology for Land '
        'Registration: A Namibian Case Study. Land, 14(8), p. 1626.',
    ]

    p = tf.paragraphs[0]
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = refs[0]
    run.font.name = 'Calibri'
    run.font.size = Pt(15)
    run.font.color.rgb = BLACK

    for ref in refs[1:]:
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = ref
        run.font.name = 'Calibri'
        run.font.size = Pt(15)
        run.font.color.rgb = BLACK


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    slide_title(prs)
    slide_overview(prs)
    slide_introduction(prs)
    slide_problem(prs)
    slide_objectives(prs)
    slide_related(prs)
    slide_solution(prs)
    slide_smart_contract(prs)
    slide_security(prs)
    slide_results(prs)
    slide_conclusion(prs)
    slide_references(prs)

    output = 'BlockLand_HIT_Presentation.pptx'
    prs.save(output)
    print(f'Saved: {output}  ({prs.slides.__len__()} slides)')

if __name__ == '__main__':
    main()
